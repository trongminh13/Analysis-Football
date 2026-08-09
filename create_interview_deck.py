from pathlib import Path
import sys

sys.path.insert(0, "/tmp/codex-slide-deps")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from PIL import Image


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "Do_Trong_Minh_Applied_AI_Interview_Deck.pptx"
MOBILE = ROOT.parent / "MOBILE-AUTO-TESTCASE-GENERATOR"

W, H = 13.333, 7.5
NAVY = "081321"
NAVY_2 = "0D1C2E"
PANEL = "12263A"
PANEL_2 = "183047"
WHITE = "F6F8FB"
MUTED = "A8B8C8"
CYAN = "42D3E8"
LIME = "B7F34A"
ORANGE = "FFB65C"
RED = "FF7272"
GRID = "264158"


def rgb(hex_value):
    return RGBColor.from_string(hex_value)


def add_bg(slide, color=NAVY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=18, color=WHITE, bold=False,
             font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
             margin=0, line_spacing=1.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rich_text(slide, runs, x, y, w, h, size=18, margin=0,
                  valign=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    for item in runs:
        run = p.add_run()
        run.text = item[0]
        run.font.name = "Aptos"
        run.font.size = Pt(item[1] if len(item) > 1 else size)
        run.font.bold = item[2] if len(item) > 2 else False
        run.font.color.rgb = rgb(item[3] if len(item) > 3 else WHITE)
    return box


def add_rect(slide, x, y, w, h, fill=PANEL, radius=True, line=None, transparency=0):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, color=GRID, width=1.5, dash=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dash:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return line


def add_title(slide, number, eyebrow, title, subtitle=None):
    add_text(slide, f"0{number}", 0.55, 0.38, 0.55, 0.32, 10, CYAN, True)
    add_line(slide, 1.15, 0.55, 1.72, 0.55, CYAN, 1.5)
    add_text(slide, eyebrow.upper(), 1.84, 0.38, 4.8, 0.32, 10, MUTED, True)
    add_text(slide, title, 0.55, 0.82, 12.15, 0.62, 28, WHITE, True)
    if subtitle:
        add_text(slide, subtitle, 0.57, 1.43, 11.8, 0.36, 11, MUTED)


def add_footer(slide, section):
    add_line(slide, 0.55, 7.1, 12.78, 7.1, GRID, 0.8)
    add_text(slide, "ĐỖ TRỌNG MINH · APPLIED AI ENGINEER", 0.55, 7.17, 4.4, 0.18, 7.5, MUTED, True)
    add_text(slide, section.upper(), 9.2, 7.17, 3.58, 0.18, 7.5, MUTED, True, align=PP_ALIGN.RIGHT)


def add_notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.text = text


def add_pill(slide, text, x, y, w, color=CYAN):
    add_rect(slide, x, y, w, 0.34, NAVY_2, True, color)
    add_text(slide, text, x, y + 0.02, w, 0.25, 9, color, True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_icon_tile(slide, num, title, detail, x, y, accent):
    add_rect(slide, x, y, 2.73, 1.3, PANEL, True, GRID)
    add_rect(slide, x + 0.18, y + 0.18, 0.5, 0.5, accent, True)
    add_text(slide, num, x + 0.18, y + 0.18, 0.5, 0.5, 11, NAVY, True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, title, x + 0.82, y + 0.17, 1.73, 0.34, 13, WHITE, True)
    add_text(slide, detail, x + 0.82, y + 0.55, 1.72, 0.52, 9.5, MUTED)


def add_cropped_picture(slide, path, x, y, w, h):
    path = Path(path)
    with Image.open(path) as im:
        iw, ih = im.size
    target = w / h
    source = iw / ih
    if source > target:
        crop = (1 - target / source) / 2
        picture = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
        picture.crop_left = crop
        picture.crop_right = crop
    else:
        crop = (1 - source / target) / 2
        picture = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
        picture.crop_top = crop
        picture.crop_bottom = crop


prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
blank = prs.slide_layouts[6]


# Slide 1
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_rect(slide, 8.55, -0.65, 5.7, 8.2, NAVY_2, False)
for i in range(7):
    add_line(slide, 8.6 + i * 0.78, 0, 8.6 + i * 0.78, 7.5, GRID, 0.55)
for i in range(10):
    add_line(slide, 8.55, i * 0.78, 13.33, i * 0.78, GRID, 0.55)
add_text(slide, "APPLIED AI · PORTFOLIO", 0.65, 0.62, 4.5, 0.3, 10, CYAN, True)
add_text(slide, "Đỗ Trọng Minh", 0.65, 1.28, 7.25, 0.72, 36, WHITE, True)
add_text(slide, "Applied AI Engineer", 0.68, 2.05, 6.5, 0.47, 20, LIME, True)
add_text(slide,
         "Building AI systems that turn unstructured input\ninto validated, editable workflows.",
         0.68, 2.82, 6.75, 1.05, 19, WHITE, False)
skills = [
    ("01", "LLM Agents", "Structured generation\n& human review", CYAN),
    ("02", "Computer Vision", "Tracking & spatial\nreasoning", LIME),
    ("03", "Backend / API", "Workflow services\n& integration", ORANGE),
    ("04", "Product AI", "Evaluation to\nreal-world use", RED),
]
for idx, item in enumerate(skills):
    num, title, detail, accent = item
    add_icon_tile(slide, num, title, detail,
                  0.68 + (idx % 2) * 3.0,
                  4.45 + (idx // 2) * 1.48,
                  accent)
add_rect(slide, 9.4, 1.15, 2.95, 4.95, PANEL, True, CYAN)
add_text(slide, "POSITIONING", 9.75, 1.55, 2.25, 0.28, 9, MUTED, True)
add_text(slide, "MODEL", 9.75, 2.19, 1.4, 0.24, 9, CYAN, True)
add_text(slide, "+", 10.68, 2.54, 0.5, 0.4, 20, MUTED, True, align=PP_ALIGN.CENTER)
add_text(slide, "REPRESENTATION", 9.75, 3.02, 2.2, 0.24, 9, LIME, True)
add_text(slide, "+", 10.68, 3.37, 0.5, 0.4, 20, MUTED, True, align=PP_ALIGN.CENTER)
add_text(slide, "VALIDATION", 9.75, 3.86, 2.2, 0.24, 9, ORANGE, True)
add_text(slide, "+", 10.68, 4.21, 0.5, 0.4, 20, MUTED, True, align=PP_ALIGN.CENTER)
add_text(slide, "SOFTWARE", 9.75, 4.7, 2.2, 0.24, 9, RED, True)
add_text(slide, "→ usable AI", 9.75, 5.33, 2.2, 0.34, 16, WHITE, True)
add_footer(slide, "Profile & positioning")
add_notes(slide, """Thời lượng mục tiêu: 60 giây.

Em có nền tảng từ Software Development và Quality Assurance, sau đó chuyển sang Applied AI. Kinh nghiệm gần nhất của em tập trung vào hệ thống multi-agent tạo test case từ tài liệu nghiệp vụ, kết hợp LLM, backend service, structured output và human review. Em cũng từng xây dựng pipeline Computer Vision cho phân tích video bóng đá, gồm detection, tracking, camera compensation và chuyển đổi tọa độ ảnh sang mặt sân thực.

Định hướng của em là phát triển các hệ thống AI có thể hiểu dữ liệu đa phương thức, thao tác trên một representation có cấu trúc và tích hợp được vào quy trình làm việc thực tế.

Thông điệp cần giữ: em không chỉ gọi model; em quan tâm representation, validation, integration và workflow của người dùng.""")


# Slide 2
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_title(slide, 2, "Case study 01", "Mobile Auto Testcase Generator",
          "From product requirements to reviewable automation artifacts")
add_rect(slide, 0.55, 2.05, 5.45, 4.62, PANEL, True, GRID)
add_text(slide, "BUSINESS PROBLEM", 0.88, 2.38, 2.8, 0.28, 10, ORANGE, True)
add_text(slide, "Manual QA authoring\nwas the bottleneck.", 0.88, 2.83, 4.5, 0.95, 25, WHITE, True)
problems = [
    ("~3 days", "for each feature workflow"),
    ("PRD → logic", "requires domain interpretation"),
    ("Structured output", "must feed downstream automation"),
]
for i, (a, b) in enumerate(problems):
    y = 4.18 + i * 0.69
    add_text(slide, a, 0.9, y, 1.52, 0.29, 12, LIME if i == 0 else CYAN, True)
    add_text(slide, b, 2.45, y, 2.95, 0.29, 10.5, MUTED)
add_rect(slide, 6.28, 2.05, 6.5, 4.62, NAVY_2, True, CYAN)
add_text(slide, "MY RESPONSIBILITY", 6.64, 2.38, 3.0, 0.28, 10, CYAN, True)
responsibilities = [
    ("01", "Agent workflow", "Three-stage reasoning pipeline"),
    ("02", "FastAPI services", "Agent execution & contracts"),
    ("03", "Structured output", "JSON-first schemas → YAML"),
    ("04", "HITL", "Review, feedback & partial rerun"),
    ("05", "Evaluation", "Internal pilot and failure analysis"),
]
for i, (n, title, detail) in enumerate(responsibilities):
    y = 2.93 + i * 0.67
    add_rect(slide, 6.65, y, 0.43, 0.43, CYAN if i < 3 else LIME, True)
    add_text(slide, n, 6.65, y, 0.43, 0.43, 8.5, NAVY, True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, title, 7.3, y - 0.01, 1.75, 0.25, 11.5, WHITE, True)
    add_text(slide, detail, 9.05, y - 0.01, 3.15, 0.25, 10, MUTED)
add_footer(slide, "Problem before technology")
add_notes(slide, """Bắt đầu bằng pain point, không đọc danh sách công nghệ.

QA phải đọc PRD, diễn giải requirement và viết test case thủ công. Với feature phức tạp, chu trình từ phân tích tới automation artifact có thể kéo dài khoảng ba ngày. Đầu ra không chỉ cần dễ đọc mà còn phải có schema đủ chặt để backend lưu trữ, reviewer chỉnh sửa và hệ thống automation chạy tiếp.

Trách nhiệm của em là thiết kế workflow ba agent, xây Agent Service bằng FastAPI, chuẩn hóa contract JSON/YAML, phát triển human-in-the-loop và partial rerun, sau đó tham gia đánh giá nội bộ.

Nếu bị hỏi stack đầy đủ: frontend React/TypeScript; backend Node.js/Express; agent service Python/FastAPI; Supabase; SSE; Langfuse. Nhưng chỉ nói khi được hỏi.""")


# Slide 3
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_title(slide, 3, "Architecture & decisions", "A staged, inspectable agent workflow",
          "Reasoning is separated from executable generation")
stages = [
    ("PRD", "INGESTION", CYAN),
    ("A1", "REQUIREMENTS", CYAN),
    ("A2", "SCENARIOS", LIME),
    ("A3", "SCRIPTS", ORANGE),
    ("✓", "VALIDATION", RED),
    ("H", "HUMAN REVIEW", WHITE),
]
start_x = 0.58
for i, (token, label, accent) in enumerate(stages):
    x = start_x + i * 2.08
    add_rect(slide, x, 2.0, 1.62, 1.15, PANEL, True, accent)
    add_text(slide, token, x, 2.16, 1.62, 0.4, 18, accent, True,
             align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.08, 2.7, 1.46, 0.22, 8, MUTED, True,
             align=PP_ALIGN.CENTER)
    if i < len(stages) - 1:
        add_line(slide, x + 1.64, 2.57, x + 2.03, 2.57, CYAN, 2)
add_text(slide, "↺ partial regeneration", 10.16, 3.25, 2.1, 0.3, 9.5, LIME, True,
         align=PP_ALIGN.CENTER)
decisions = [
    ("01", "Segment the context", "Process flows independently instead of one oversized prompt."),
    ("02", "JSON first", "Validate semantic structure before rendering executable YAML."),
    ("03", "Rerun by stage", "Preserve reviewed artifacts; regenerate only stale downstream work."),
]
for i, (n, title, detail) in enumerate(decisions):
    x = 0.58 + i * 4.14
    add_rect(slide, x, 3.83, 3.78, 1.4, NAVY_2, True, GRID)
    add_text(slide, n, x + 0.22, 4.07, 0.42, 0.28, 10, CYAN, True)
    add_text(slide, title, x + 0.72, 4.03, 2.6, 0.31, 13, WHITE, True)
    add_text(slide, detail, x + 0.22, 4.5, 3.3, 0.5, 9.5, MUTED)
add_rect(slide, 0.58, 5.55, 12.18, 1.08, PANEL, True, GRID)
metrics = [
    ("~3 days", "BEFORE", MUTED),
    ("< 5 min", "PIPELINE OUTPUT", LIME),
    ("88.4%", "FIRST-RUN YAML*", CYAN),
    ("28", "BETA USERS", ORANGE),
]
for i, (v, l, c) in enumerate(metrics):
    x = 0.95 + i * 3.03
    if i:
        add_line(slide, x - 0.35, 5.77, x - 0.35, 6.39, GRID, 1)
    add_text(slide, v, x, 5.71, 2.1, 0.37, 20, c, True)
    add_text(slide, l, x, 6.13, 2.15, 0.2, 8, MUTED, True)
add_text(slide, "*Reported beta metric; sample count must be confirmed before interview.",
         7.26, 6.7, 5.5, 0.18, 7.5, MUTED, False, align=PP_ALIGN.RIGHT)
add_footer(slide, "Architecture & evidence")
add_notes(slide, """Ba quyết định cần nhấn mạnh:
1) Segmented context: tài liệu dài được chia theo flow để tránh token overflow và cô lập lỗi.
2) JSON-first: Agent 2/3 tạo representation có schema; validate trước khi render YAML. Format đúng không đồng nghĩa logic đúng, nên vẫn cần human review.
3) Partial rerun: version artifact theo stage, chỉ làm stale downstream thay vì chạy lại toàn bộ.

Kết quả trình bày trên slide:
- Workflow thủ công khoảng ba ngày xuống dưới năm phút cho pipeline output.
- Pilot 28 QA Engineers/Developers trong báo cáo nội bộ.
- 88,4% là tỷ lệ YAML chạy thành công ngay trên máy ảo theo báo cáo beta; KHÔNG gọi là accuracy tổng thể.

CÂU HỎI: “88,4% được định nghĩa và đo như thế nào?”

Trả lời trung thực:
“Đơn vị đánh giá là từng file YAML được hệ thống sinh trong beta. Một mẫu pass khi chạy thành công ngay trên máy ảo mà không cần chỉnh sửa. Người dùng beta là 28 QA Engineers và Developers, kết quả được tổng hợp trong báo cáo nội bộ. Các lỗi phổ biến của 11,6% còn lại là sai element ID hoặc thiếu assertion step. Repo hiện lưu tỷ lệ 88,4% nhưng chưa lưu tổng số YAML được chấm, vì vậy trước buổi phỏng vấn em cần xác nhận mẫu số N với nhóm; em không muốn biến một tỷ lệ pilot thành tuyên bố thống kê mạnh hơn bằng chứng. Metric cũng chỉ đo khả năng chạy lần đầu, chưa chứng minh semantic correctness, coverage đầy đủ hay generalization trên một test set độc lập.”

MẪU SỐ: [ĐIỀN N YAML SAU KHI XÁC NHẬN].

Nếu interviewer hỏi zero-shot: “Tài liệu hiện có không tách riêng zero-shot benchmark. Em sẽ gọi đúng là first-run execution rate trong beta; muốn đánh giá zero-shot cần cố định prompt/model/test set và không cho feedback trước lần chạy.”""")


# Slide 4
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_title(slide, 4, "Case study 02", "Football Video Analysis",
          "Unstructured visual input → coordinates, trajectories and events")
pipeline = [
    ("VIDEO", "RAW INPUT", CYAN),
    ("YOLO", "DETECTION", CYAN),
    ("ID", "TRACKING", LIME),
    ("ΔXY", "CAMERA", ORANGE),
    ("H", "HOMOGRAPHY", RED),
    ("m", "REAL WORLD", WHITE),
]
for i, (big, small, accent) in enumerate(pipeline):
    x = 0.55 + i * 2.1
    add_rect(slide, x, 2.02, 1.62, 1.02, PANEL, True, accent)
    add_text(slide, big, x, 2.15, 1.62, 0.35, 16, accent, True, align=PP_ALIGN.CENTER)
    add_text(slide, small, x, 2.61, 1.62, 0.2, 7.5, MUTED, True, align=PP_ALIGN.CENTER)
    if i < 5:
        add_line(slide, x + 1.65, 2.52, x + 2.04, 2.52, GRID, 2)
add_rect(slide, 0.55, 3.42, 5.15, 2.92, NAVY_2, True, GRID)
image_path = ROOT / "training/runs/detect/train7/val_batch1_pred.jpg"
add_cropped_picture(slide, image_path, 0.74, 3.62, 4.77, 2.52)
add_rect(slide, 5.99, 3.42, 6.79, 2.92, PANEL, True, GRID)
add_text(slide, "TRANSFERABLE FOUNDATION", 6.35, 3.75, 3.8, 0.25, 9.5, LIME, True)
foundations = [
    ("Detection", "objects from pixels"),
    ("Tracking", "identity through time"),
    ("Spatial transform", "image plane → ground plane"),
    ("Coordinates", "consistent representations"),
    ("Error propagation", "upstream noise affects metrics"),
]
for i, (a, b) in enumerate(foundations):
    y = 4.2 + i * 0.39
    add_text(slide, "•", 6.34, y, 0.22, 0.22, 11, CYAN, True)
    add_text(slide, a, 6.61, y, 1.65, 0.22, 10.5, WHITE, True)
    add_text(slide, b, 8.28, y, 3.78, 0.22, 9.7, MUTED)
add_text(slide, "AEC connection", 6.35, 6.24, 1.3, 0.22, 8.5, ORANGE, True)
add_text(slide, "Visual data → geometry / structured representation",
         7.62, 6.23, 4.65, 0.24, 10, WHITE, True)
add_footer(slide, "Spatial reasoning foundation")
add_notes(slide, """Luồng: video → YOLO detection → ByteTrack → camera motion compensation bằng Lucas-Kanade Optical Flow → homography → real-world coordinates → speed, distance và possession.

Thông điệp chính:
“Project này giúp em làm việc với một pipeline mà đầu vào là dữ liệu thị giác không có cấu trúc, nhưng đầu ra cần trở thành tọa độ và sự kiện có ý nghĩa. Tư duy đó gần với các bài toán chuyển ảnh hoặc bản vẽ thành geometry hay representation có cấu trúc trong AEC.”

Không nói project tương đương 2D-to-3D trong AEC. Chỉ nói đây là nền tảng về detection, tracking, spatial transformation, coordinate systems và error propagation.

Nếu hỏi error propagation: detection jitter ảnh hưởng foot point; camera estimate sai ảnh hưởng adjusted coordinate; homography/calibration sai ảnh hưởng mét; FPS sai ảnh hưởng speed. Vì vậy mAP detector tốt chưa chứng minh tốc độ chính xác.

Chi tiết trung thực: artifact huấn luyện ghi yolov5x.pt, 4 class; README ghi YOLOv8 và cần được đồng bộ. Video mẫu 25 FPS nhưng code hiện hard-code 24 FPS.""")


# Slide 5
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_title(slide, 5, "AEC problem-solving framework", "Representation before model selection",
          "A measurable path from domain need to editable output")
steps = [
    ("01", "Problem", CYAN), ("02", "Representation", LIME),
    ("03", "Baseline", ORANGE), ("04", "Architecture", RED),
    ("05", "Evaluation", CYAN), ("06", "Expert loop", LIME),
]
for i, (n, label, accent) in enumerate(steps):
    x = 0.58 + i * 2.08
    add_rect(slide, x, 1.94, 1.64, 0.73, PANEL, True, accent)
    add_text(slide, n, x + 0.11, 2.13, 0.37, 0.22, 9, accent, True)
    add_text(slide, label, x + 0.5, 2.09, 1.0, 0.26, 10, WHITE, True)
add_rect(slide, 0.58, 2.97, 3.72, 3.69, PANEL, True, GRID)
add_text(slide, "DEFINE THE DECISION", 0.9, 3.28, 2.6, 0.25, 9.5, CYAN, True)
questions = [
    "Who is the end user?",
    "What decision does AI support?",
    "What is the real input?",
    "Must output remain editable?",
    "Which constraints are mandatory?",
    "Which errors are unacceptable?",
]
for i, q in enumerate(questions):
    y = 3.72 + i * 0.43
    add_text(slide, f"{i+1:02}", 0.9, y, 0.34, 0.22, 8, MUTED, True)
    add_text(slide, q, 1.37, y, 2.5, 0.25, 10.5, WHITE)
add_rect(slide, 4.57, 2.97, 4.08, 3.69, NAVY_2, True, LIME)
add_text(slide, "DOMAIN REPRESENTATION", 4.91, 3.28, 3.2, 0.25, 9.5, LIME, True)
repr_items = ["Image / CAD / BIM", "objects", "geometry", "attributes",
              "spatial relationships", "constraints"]
for i, item in enumerate(repr_items):
    y = 3.73 + i * 0.43
    add_rect(slide, 4.95, y - 0.04, 3.26, 0.31, PANEL_2, True)
    add_text(slide, item, 5.05, y, 3.05, 0.2, 9.5,
             WHITE if i == 0 else MUTED, i == 0, align=PP_ALIGN.CENTER)
    if i < len(repr_items) - 1:
        add_text(slide, "↓", 6.43, y + 0.29, 0.3, 0.18, 9, CYAN, True,
                 align=PP_ALIGN.CENTER)
add_rect(slide, 8.92, 2.97, 3.86, 3.69, PANEL, True, ORANGE)
add_text(slide, "MEASURABLE BASELINE", 9.25, 3.28, 3.0, 0.25, 9.5, ORANGE, True)
baseline = ["Input parser", "Structured IR", "Rule generator",
            "Constraint validator", "Editable output"]
for i, item in enumerate(baseline):
    y = 3.78 + i * 0.51
    add_rect(slide, 9.26, y - 0.07, 0.37, 0.37, ORANGE if i == 0 else NAVY_2, True,
             ORANGE)
    add_text(slide, str(i + 1), 9.26, y - 0.07, 0.37, 0.37, 8, NAVY if i == 0 else ORANGE,
             True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, item, 9.82, y, 2.4, 0.22, 10.5, WHITE, i in (1, 3))
add_text(slide, "Model complexity follows measured gaps — not fashion.",
         9.26, 6.23, 3.0, 0.23, 8.5, MUTED)
add_footer(slide, "How I would enter AEC")
add_notes(slide, """Đây là slide quan trọng nhất. Nói chậm.

1. Problem definition: xác định end user, quyết định AI hỗ trợ, input thật, output editable, constraint bắt buộc và mức lỗi không chấp nhận.
2. Domain representation: trước khi chọn model, xác định representation phù hợp. Với AEC, output thường không nên chỉ là ảnh đẹp mà phải là dữ liệu kiểm tra được: room graph, polygon, wall/opening, adjacency, dimension, semantic object và design constraint.
3. Baseline: không bắt đầu ngay bằng diffusion hoặc fine-tuning. Xây baseline nhỏ nhất có thể đo: parser → structured intermediate representation → rule-based generator → constraint validator → editable output.
4. Prototype architecture: tách parser/model/validator/storage/UI để thay model mà không phá workflow.
5. Evaluation: metric theo nghiệp vụ; geometry validity, constraint violations, field accuracy, edit time, false negative nghiêm trọng.
6. Expert loop: architect/engineer review, lưu correction, active learning và regression set.

Ví dụ nếu input là bản vẽ PDF: trước tiên xác định PDF vector hay scan. Vector thì ưu tiên parse text/path/layer; scan mới cần OCR/detection/segmentation. Với bản vẽ rất lớn dùng tiled inference và merge về global coordinates. Representation có thể là graph với room/wall/opening nodes và adjacency edges.""")


# Slide 6
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_title(slide, 6, "Role fit & execution", "Why this role — and my first 90 days",
          "Join through the workflow; earn model complexity with evidence")
add_rect(slide, 0.58, 2.02, 4.05, 4.55, PANEL, True, CYAN)
add_text(slide, "WHY THIS ROLE", 0.93, 2.37, 2.5, 0.25, 10, CYAN, True)
add_text(slide, "AI connected to\ndomain tools — not\njust conversation.",
         0.93, 2.85, 3.12, 1.38, 23, WHITE, True)
add_text(slide, "Image-to-script · agentic workflows\ngenerative geometry · validation\nsoftware engineering",
         0.93, 4.56, 3.03, 0.95, 11, MUTED)
add_pill(slide, "MODEL", 0.93, 5.86, 0.75, CYAN)
add_pill(slide, "DATA", 1.81, 5.86, 0.72, LIME)
add_pill(slide, "RULES", 2.66, 5.86, 0.78, ORANGE)
add_pill(slide, "UX", 3.57, 5.86, 0.55, RED)
phases = [
    ("0—30", "UNDERSTAND", CYAN,
     ["Architect workflow", "Current data & tools", "One valuable use case", "Initial evaluation set"]),
    ("31—60", "BASELINE", LIME,
     ["End-to-end baseline", "Logging & versioning", "Metrics", "Internal feedback"]),
    ("61—90", "VALIDATE", ORANGE,
     ["Compare model paths", "Accuracy & latency", "Editability", "Prototype → tool roadmap"]),
]
for i, (days, label, accent, items) in enumerate(phases):
    x = 4.94 + i * 2.61
    add_rect(slide, x, 2.02, 2.35, 4.55, NAVY_2, True, accent)
    add_text(slide, days, x + 0.25, 2.33, 1.4, 0.42, 21, accent, True)
    add_text(slide, "DAYS", x + 1.54, 2.48, 0.55, 0.19, 8, MUTED, True)
    add_text(slide, label, x + 0.25, 2.98, 1.85, 0.26, 10, WHITE, True)
    add_line(slide, x + 0.25, 3.38, x + 2.08, 3.38, GRID, 1)
    for j, item in enumerate(items):
        y = 3.68 + j * 0.57
        add_text(slide, "•", x + 0.25, y, 0.2, 0.22, 11, accent, True)
        add_text(slide, item, x + 0.52, y, 1.54, 0.38, 9.7, WHITE)
add_text(slide, "Outcome", 9.45, 6.73, 0.65, 0.18, 8, MUTED, True)
add_text(slide, "A measured roadmap from prototype to internal product",
         10.12, 6.7, 2.66, 0.23, 8.5, WHITE, True, align=PP_ALIGN.RIGHT)
add_footer(slide, "90-day product mindset")
add_notes(slide, """Vì sao vị trí phù hợp:
“Điểm em quan tâm không chỉ là dùng LLM để trò chuyện, mà là kết nối AI với dữ liệu và công cụ chuyên ngành. Các bài toán image-to-script, agentic workflow và generative geometry đều cần kết hợp model, representation, validation và software engineering. Đây là hướng em muốn phát triển lâu dài.”

30 ngày đầu: shadow workflow kiến trúc sư/kỹ sư; hiểu input, tool, vocabulary và failure cost; chọn một use case nhỏ; tạo evaluation set với chuyên gia.

30 ngày tiếp: baseline end-to-end; logging/model-data versioning; metric kỹ thuật và metric edit-time; review định kỳ với người dùng nội bộ.

30 ngày cuối: so sánh model/rule/hybrid paths; tối ưu accuracy, latency và editability; viết roadmap với risk, data need và tiêu chí production.

Câu chốt:
“Mục tiêu 90 ngày của em không phải demo nhiều model nhất, mà là chứng minh một workflow nhỏ có giá trị, đo được và có đường phát triển thành internal tool.”""")


prs.core_properties.title = "Đỗ Trọng Minh — Applied AI Engineer"
prs.core_properties.subject = "Interview portfolio: agentic systems, computer vision and AEC problem solving"
prs.core_properties.author = "Đỗ Trọng Minh"
prs.core_properties.keywords = "Applied AI, Computer Vision, LLM Agents, AEC"
prs.save(OUT)
print(OUT)
